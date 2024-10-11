import streamlit as st
from langchain.document_loaders import UnstructuredExcelLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OllamaEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain.llms import Ollama
from langchain.prompts import PromptTemplate  
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import time

def process_uploaded_file(uploaded_file):
    if uploaded_file is not None:
        if uploaded_file.type == "application/pdf":
            reader = PdfReader(uploaded_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            with open('output.txt', 'w', encoding="utf-8") as f:
                f.write(text)
            file_loader = TextLoader('./output.txt', encoding='utf8')
            docs = file_loader.load()
            
            return docs
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            file_loader = UnstructuredExcelLoader(uploaded_file)
            docs = file_loader.load()
            return docs
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(uploaded_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            with open('output.txt', 'w') as f:
                f.write(text)
            file_loader = TextLoader('./output.txt')
            docs = file_loader.load()
            return docs
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = Presentation(uploaded_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            with open('output.txt', 'w') as f:
                f.write(text)
            file_loader = TextLoader('./output.txt')
            docs = file_loader.load()
            return docs
        else:
            st.error("Unsupported file type. Please upload a PDF, Excel (.xlsx), Word (.docx), PowerPoint (.pptx), or text file.")
            return None
    else:
        return None

def chunk_data(docs, chunk_size=800, chunk_overlap=50):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return text_splitter.split_documents(docs)

def create_vectorstore(docs):
    oembed = OllamaEmbeddings(
        base_url="http://localhost:11434",  
        model="nomic-embed-text"
    )
    documents = chunk_data(docs)
    vectorstore = FAISS.from_documents(documents, oembed)
    return vectorstore

def create_qa_chain(vectorstore):
    llm = Ollama(model="llama3")  
    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
    qa = RetrievalQA.from_chain_type(
        llm=llm, 
        chain_type="stuff", 
        retriever=retriever
    )
    return qa

def respond(question, history, qa):
    start_time = time.time()
    response = qa(question)["result"]
    end_time = time.time()
    time_taken = end_time - start_time
    return response, time_taken

st.title("Chatbot")

if 'history' not in st.session_state:
    st.session_state.history = []

if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = None

with st.sidebar:
    st.header("Conversation History")
    for i, (q, a, t) in enumerate(st.session_state.history):
        st.markdown(f"{i+1}. **User:** {q}\n   **Chatbot:** {a}\n   **Time:** {t:.2f} seconds")

with st.form("chatbot"):
    uploaded_file = st.file_uploader("Upload a PDF, Excel (.xlsx), Word (.docx), or PowerPoint (.pptx) file", type=["pdf", "xlsx", "docx", "pptx"])
    question_user = st.text_input("Ask me a question", "")

    if st.form_submit_button("Submit"):
        docs = process_uploaded_file(uploaded_file)
        if docs:
            st.session_state.vectorstore = create_vectorstore(docs)
            qa = create_qa_chain(st.session_state.vectorstore)
            response, time_taken = respond(question_user, st.session_state.history, qa)
            st.text_area("Chatbot Response", value=response, height=len(response.splitlines()) * 20)
            st.session_state.history.append((question_user, response, time_taken))  
        elif st.session_state.vectorstore:
            qa = create_qa_chain(st.session_state.vectorstore)
            response, time_taken = respond(question_user, st.session_state.history, qa)
            st.text_area("Chatbot Response", value=response, height=len(response.splitlines()) * 20)
            st.session_state.history.append((question_user, response, time_taken))  
        else:
            st.error("Please upload a valid PDF, Excel (.xlsx), Word (.docx), PowerPoint (.pptx), or text file.")